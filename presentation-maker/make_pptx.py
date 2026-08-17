import json, sys, os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

TEAL_DARK = RGBColor(0x0B, 0x3C, 0x46)
TEAL = RGBColor(0x12, 0x7A, 0x8D)
TEAL2 = RGBColor(0x3C, 0x8F, 0x96)
ORANGE = RGBColor(0xF0, 0x7C, 0x1D)
DARK = RGBColor(0x1F, 0x2A, 0x30)
GRAY = RGBColor(0x6B, 0x7A, 0x80)
LIGHT = RGBColor(0xE2, 0xEF, 0xF4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Segoe UI"
FONT_LIGHT = "Segoe UI Light"
LOGO = os.path.expanduser("~/presentations/logo.png")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height
PAGE = [0]

def set_alpha(shape, opacity):
    spPr = shape._element.spPr
    sf = spPr.find(qn('a:solidFill'))
    if sf is None:
        return
    clr = sf.find(qn('a:srgbClr'))
    if clr is None:
        return
    old_a = clr.find(qn('a:alpha'))
    if old_a is not None:
        clr.remove(old_a)
    a = clr.makeelement(qn('a:alpha'), {'val': str(int(opacity * 1000))})
    clr.append(a)

def no_line(shape):
    shape.line.fill.background()

def gradient_bg(slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    no_line(bg)
    f = bg.fill
    f.gradient()
    f.gradient_angle = 135
    f.gradient_stops[0].color.rgb = TEAL_DARK
    f.gradient_stops[1].color.rgb = TEAL

def deco(slide, on_dark):
    c1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.4), Inches(-3.6), Inches(8.8), Inches(8.8))
    no_line(c1)
    c1.fill.solid()
    c1.fill.fore_color.rgb = WHITE if on_dark else TEAL2
    set_alpha(c1, 8 if on_dark else 10)
    c2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.9), Inches(5.3), Inches(2.3), Inches(2.3))
    no_line(c2)
    c2.fill.solid()
    c2.fill.fore_color.rgb = ORANGE
    set_alpha(c2, 55 if on_dark else 45)

def textbox(slide, l, t, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT, font=FONT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return tf

def footer(slide, light=False):
    PAGE[0] += 1
    color = WHITE if light else GRAY
    textbox(slide, 0.6, 7.0, 6, 0.4, "1c-gendalf.ru", 10, color)
    tb = slide.shapes.add_textbox(Inches(11.9), Inches(7.0), Inches(0.9), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    p.text = str(PAGE[0])
    p.font.size = Pt(10)
    p.font.color.rgb = color
    p.font.name = FONT
    p.alignment = PP_ALIGN.RIGHT

def add_logo(slide, big=False):
    if os.path.exists(LOGO):
        if big:
            slide.shapes.add_picture(LOGO, Inches(0.6), Inches(0.5), height=Inches(1.1))
        else:
            slide.shapes.add_picture(LOGO, Inches(11.7), Inches(0.3), height=Inches(0.55))

def title_slide(d):
    s = prs.slides.add_slide(BLANK)
    gradient_bg(s)
    deco(s, True)
    add_logo(s, big=True)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(2.4), Pt(9), Inches(1.8))
    no_line(bar)
    bar.fill.solid()
    bar.fill.fore_color.rgb = ORANGE
    textbox(s, 1.0, 2.35, 11.5, 1.9, d.get("title", ""), 50, WHITE, bold=True)
    textbox(s, 1.0, 4.5, 11.5, 0.8, d.get("subtitle", ""), 22, LIGHT, font=FONT_LIGHT)
    meta = "   |   ".join(x for x in [d.get("author", ""), d.get("date", "")] if x)
    if meta:
        textbox(s, 1.0, 6.5, 11, 0.5, meta, 13, TEAL2)

def section_slide(d):
    s = prs.slides.add_slide(BLANK)
    gradient_bg(s)
    deco(s, True)
    num = d.get("num", "")
    if num:
        textbox(s, 0.6, 0.9, 6, 2.8, num, 120, TEAL2, bold=True)
    textbox(s, 0.6, 3.7, 12, 1.5, d.get("title", ""), 40, WHITE, bold=True)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(5.3), Inches(2.4), Pt(7))
    no_line(bar)
    bar.fill.solid()
    bar.fill.fore_color.rgb = ORANGE
    footer(s, light=True)

def heading(slide, text):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.55), Pt(8), Inches(0.85))
    no_line(bar)
    bar.fill.solid()
    bar.fill.fore_color.rgb = ORANGE
    textbox(slide, 0.95, 0.5, 11.5, 1.0, text, 34, DARK, bold=True)

def bullets_slide(d):
    s = prs.slides.add_slide(BLANK)
    deco(s, False)
    add_logo(s)
    heading(s, d.get("title", ""))
    tb = s.shapes.add_textbox(Inches(0.95), Inches(1.9), Inches(11.2), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(d.get("items", [])):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r1 = p.add_run()
        r1.text = "•  "
        r1.font.color.rgb = ORANGE
        r1.font.size = Pt(22)
        r1.font.bold = True
        r1.font.name = FONT
        r2 = p.add_run()
        r2.text = item
        r2.font.color.rgb = DARK
        r2.font.size = Pt(20)
        r2.font.name = FONT
        p.space_after = Pt(16)
    footer(s)

def stats_slide(d):
    s = prs.slides.add_slide(BLANK)
    deco(s, False)
    add_logo(s)
    heading(s, d.get("title", ""))
    stats = d.get("stats", [])
    n = max(len(stats), 1)
    width = 12.1 / n
    for i, st in enumerate(stats):
        x = 0.6 + i * width
        textbox(s, x, 2.5, width - 0.5, 1.7, st.get("value", ""), 72, ORANGE if i % 2 == 0 else TEAL, bold=True)
        textbox(s, x, 4.3, width - 0.5, 1.4, st.get("label", ""), 16, GRAY, font=FONT_LIGHT)
    footer(s)

def bento_slide(d):
    s = prs.slides.add_slide(BLANK)
    add_logo(s)
    heading(s, d.get("title", ""))
    cards = d.get("cards", [])
    pos = [(0.6, 1.8), (6.85, 1.8), (0.6, 4.55), (6.85, 4.55)]
    for i, c in enumerate(cards[:4]):
        x, y = pos[i]
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(5.9), Inches(2.5))
        no_line(card)
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT
        textbox(s, x + 0.35, y + 0.3, 5.2, 0.7, c.get("title", ""), 22, TEAL, bold=True)
        textbox(s, x + 0.35, y + 1.1, 5.2, 1.2, c.get("text", ""), 15, DARK)
    footer(s)

def image_slide(d):
    s = prs.slides.add_slide(BLANK)
    if os.path.exists(d.get("image", "")):
        s.shapes.add_picture(d["image"], 0, 0, SW, SH)
    ov = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.3), SW, Inches(3.2))
    no_line(ov)
    ov.fill.solid()
    ov.fill.fore_color.rgb = TEAL_DARK
    set_alpha(ov, 75)
    textbox(s, 0.6, 4.9, 12, 1.5, d.get("title", ""), 40, WHITE, bold=True)
    if d.get("caption"):
        textbox(s, 0.6, 6.3, 12, 0.6, d["caption"], 16, LIGHT, font=FONT_LIGHT)

def table_slide(d):
    s = prs.slides.add_slide(BLANK)
    deco(s, False)
    add_logo(s)
    heading(s, d.get("title", ""))
    header = d.get("header", [])
    rows = d.get("rows", [])
    ts = s.shapes.add_table(len(rows) + 1, len(header), Inches(0.6), Inches(1.9), Inches(12.1), Inches(0.55 * (len(rows) + 1)))
    tbl = ts.table
    tbl.first_row = False
    tbl.horz_banding = False
    for c, h in enumerate(header):
        cell = tbl.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = TEAL_DARK
        for p in cell.text_frame.paragraphs:
            p.font.color.rgb = WHITE
            p.font.bold = True
            p.font.size = Pt(16)
            p.font.name = FONT
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = tbl.cell(r + 1, c)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 == 0 else LIGHT
            for p in cell.text_frame.paragraphs:
                p.font.color.rgb = DARK
                p.font.size = Pt(15)
                p.font.name = FONT
    footer(s)

def chart_slide(d):
    s = prs.slides.add_slide(BLANK)
    deco(s, False)
    add_logo(s)
    heading(s, d.get("title", ""))
    if os.path.exists(d.get("image", "")):
        s.shapes.add_picture(d["image"], Inches(1.0), Inches(1.8), width=Inches(11.3))
    footer(s)

def final_slide(d):
    s = prs.slides.add_slide(BLANK)
    gradient_bg(s)
    deco(s, True)
    add_logo(s, big=True)
    textbox(s, 0.6, 2.6, 12, 1.5, d.get("title", "Спасибо за внимание!"), 44, WHITE, bold=True)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(4.2), Inches(2.4), Pt(8))
    no_line(bar)
    bar.fill.solid()
    bar.fill.fore_color.rgb = ORANGE
    if d.get("contacts"):
        textbox(s, 0.6, 4.6, 12, 1.0, d["contacts"], 18, LIGHT, font=FONT_LIGHT)

BUILDERS = {
    "title": title_slide, "section": section_slide, "bullets": bullets_slide,
    "table": table_slide, "chart": chart_slide, "final": final_slide,
    "stats": stats_slide, "bento": bento_slide, "image": image_slide,
}

def main():
    if len(sys.argv) < 3:
        print("Usage: make_pptx.py <input.json> <output.pptx>")
        sys.exit(1)
    with open(sys.argv[1], encoding="utf-8-sig") as f:
        data = json.load(f)
    n = 0
    for sl in data.get("slides", []):
        if sl.get("type") == "section":
            n += 1
            sl.setdefault("num", f"{n:02d}")
    title_slide(data)
    for sl in data.get("slides", []):
        b = BUILDERS.get(sl.get("type"))
        if b:
            b(sl)
    if data.get("final"):
        final_slide(data["final"])
    prs.save(sys.argv[2])
    print(f"Saved: {sys.argv[2]}")

main()